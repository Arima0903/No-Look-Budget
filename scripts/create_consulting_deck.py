#!/usr/bin/env python3
"""
Orbit Budget コンサルティング資料 PowerPoint 生成スクリプト
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# カラーパレット
NAVY = RGBColor(8, 11, 20)
NAVY_LIGHT = RGBColor(26, 26, 46)
GREEN = RGBColor(74, 222, 128)
GREEN_DARK = RGBColor(34, 197, 94)
RED = RGBColor(239, 68, 68)
ORANGE = RGBColor(251, 146, 60)
BLUE = RGBColor(96, 165, 250)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(156, 163, 175)
DARK_GRAY = RGBColor(55, 65, 81)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=NAVY):
    """スライド背景を設定"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    """テキストボックスを追加"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None, radius=0.1):
    """角丸矩形を追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, content_lines,
             title_color=GREEN, accent_color=None):
    """カード型コンポーネント"""
    add_shape_rect(slide, left, top, width, height, NAVY_LIGHT, RGBColor(40, 40, 60))
    add_text(slide, left + 0.2, top + 0.15, width - 0.4, 0.4, title,
             font_size=14, color=title_color, bold=True)
    y = top + 0.55
    for line in content_lines:
        color = accent_color if accent_color else GRAY
        if line.startswith("!"):
            line = line[1:]
            color = WHITE
        add_text(slide, left + 0.2, y, width - 0.4, 0.3, line,
                 font_size=11, color=color)
        y += 0.28
    return y


# ==========================================
# スライド1: タイトル
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, 0, 2.5, 13.333, 3.0, RGBColor(15, 18, 30))
add_text(slide, 0, 0.8, 13.333, 0.5, "STRATEGIC REVIEW", font_size=14,
         color=GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, 0, 2.7, 13.333, 1.0,
         "「浪費家の浪費を止める」は実現できるのか？", font_size=36,
         color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, 0, 3.8, 13.333, 0.6,
         "Orbit Budget のプロダクト戦略 再検証", font_size=20,
         color=GREEN, alignment=PP_ALIGN.CENTER)
add_text(slide, 0, 4.8, 13.333, 0.4,
         "行動経済学 × ビジネス戦略 × ADHD当事者 の3視点による統合分析",
         font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, 0, 6.5, 13.333, 0.4, "2026.03.19 | Confidential",
         font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# ==========================================
# スライド2: Executive Summary
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.6, 0.4, 12, 0.5, "EXECUTIVE SUMMARY", font_size=28,
         color=WHITE, bold=True)
add_shape_rect(slide, 0.6, 1.0, 12.1, 0.03, GREEN)

add_text(slide, 0.6, 1.3, 12, 0.5,
         "結論: Orbit Budget は「記録を楽にするアプリ」であり、まだ「浪費を止めるアプリ」にはなっていない",
         font_size=16, color=ORANGE, bold=True)

add_card(slide, 0.6, 2.0, 3.8, 2.5, "✓ 強み（Keep）",
         ["!No-Look Experience は唯一無二",
          "!ウィジェット＝「支払いの痛み」の再注入",
          "!立替分離＝管理崩壊の防止",
          "!3秒入力＝行動コストの極限削減",
          "!情報フィードバック欠如の完全カバー"], GREEN)

add_card(slide, 4.7, 2.0, 3.8, 2.5, "△ 課題（Problem）",
         ["!購入「前」の介入メカニズムがない",
          "!日割り予算の概念がない",
          "!通知・リマインダーがない",
          "!入力忘れ→データ信頼性崩壊リスク",
          "!プロダクト成功＝ユーザー離脱の矛盾"], RED)

add_card(slide, 8.8, 2.0, 3.8, 2.5, "→ 提言（Action）",
         ["!①「今日使える金額」の導入（工数: S）",
          "!②購入前チェックモード（工数: M）",
          "!③記録リマインダー通知（工数: S）",
          " ",
          "!上記3つで「記録型→介入型」へ転換"], BLUE)

add_text(slide, 0.6, 4.8, 12, 1.0,
         "現在の Orbit Budget は浪費の16原因のうち約半数に有効。"
         "特に「情報フィードバックの欠如」は全原因をカバー。\n"
         "残り半数（衝動性の抑制・購入前介入）を小工数の3機能で補完すれば、"
         "「浪費を本当に止めるアプリ」として競合と圧倒的な差別化が可能。",
         font_size=13, color=GRAY)

# ==========================================
# スライド3: 浪費の根本原因 MECE分解
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.6, 0.4, 12, 0.5, "浪費の根本原因ツリー（MECE分解）",
         font_size=28, color=WHITE, bold=True)
add_shape_rect(slide, 0.6, 1.0, 12.1, 0.03, GREEN)
add_text(slide, 0.6, 1.2, 12, 0.4,
         "なぜ人は浪費するのか？ — 4軸 × 16原因に分解し、Orbit Budget の守備範囲を特定",
         font_size=13, color=GRAY)

categories = [
    ("A. 認知バイアス", BLUE, [
        ("A-1 現在バイアス", "○"),
        ("A-2 メンタルアカウンティング", "○"),
        ("A-3 アンカリング効果", "×"),
        ("A-4 楽観性バイアス", "○"),
        ("A-5 支払いの痛みの鈍麻", "◎"),
    ]),
    ("B. 自己制御の失敗", ORANGE, [
        ("B-1 意志力の枯渇", "△"),
        ("B-2 衝動性（ADHD）", "○"),
        ("B-3 社会的圧力", "△"),
        ("B-4 習慣化された浪費", "○"),
    ]),
    ("C. 情報の欠如", GREEN, [
        ("C-1 支出の不可視性", "◎"),
        ("C-2 記録の挫折", "◎"),
        ("C-3 予算の不在", "◎"),
        ("C-4 立替による管理崩壊", "◎"),
    ]),
    ("D. 環境・構造要因", RED, [
        ("D-1 ダークパターン", "×"),
        ("D-2 サブスク累積", "△"),
        ("D-3 収支の時間的分離", "◎"),
    ]),
]

x_start = 0.6
for i, (cat_name, color, items) in enumerate(categories):
    x = x_start + i * 3.1
    add_shape_rect(slide, x, 1.7, 2.9, 0.45, NAVY_LIGHT, color)
    add_text(slide, x + 0.1, 1.75, 2.7, 0.35, cat_name,
             font_size=13, color=color, bold=True)

    for j, (item, rating) in enumerate(items):
        y = 2.3 + j * 0.65
        bg = NAVY_LIGHT
        rating_color = GREEN if rating in ("◎", "○") else (ORANGE if rating == "△" else RED)
        add_shape_rect(slide, x, y, 2.9, 0.55, bg, RGBColor(40, 40, 60))
        add_text(slide, x + 0.15, y + 0.05, 2.1, 0.25, item,
                 font_size=10, color=WHITE)
        add_text(slide, x + 2.3, y + 0.05, 0.5, 0.25, rating,
                 font_size=14, color=rating_color, bold=True, alignment=PP_ALIGN.CENTER)
        label = {"◎": "有効", "○": "部分有効", "△": "効果薄", "×": "未対応"}[rating]
        add_text(slide, x + 2.1, y + 0.3, 0.7, 0.2, label,
                 font_size=8, color=rating_color, alignment=PP_ALIGN.CENTER)

add_text(slide, 0.6, 6.3, 12, 0.6,
         "凡例: ◎ = 有効（解決できる） / ○ = 部分的に有効 / △ = 効果薄い / × = 対応できていない\n"
         "Orbit Budget は C（情報の欠如）を完全カバー。A-5（支払いの痛みの再注入）が最大の武器。",
         font_size=11, color=GRAY)

# ==========================================
# スライド4: 競合比較 - 記録型 vs 介入型
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.6, 0.4, 12, 0.5, "競合比較: 「記録型」 vs 「介入型」",
         font_size=28, color=WHITE, bold=True)
add_shape_rect(slide, 0.6, 1.0, 12.1, 0.03, GREEN)

# マネーフォワード
add_shape_rect(slide, 0.6, 1.4, 5.8, 2.8, NAVY_LIGHT, GRAY)
add_text(slide, 0.8, 1.5, 5.4, 0.4, "マネーフォワード / Zaim（記録型）",
         font_size=16, color=GRAY, bold=True)
add_text(slide, 0.8, 2.0, 5.4, 0.3, "支出発生 → 自動記録 → 月末レポート → 反省 → 翌月も同じ",
         font_size=11, color=GRAY)
flow_items = [
    ("Motivation", "正確に記録したい（外発的）"),
    ("Ability", "自動連携で楽。だが分析にスキル必要"),
    ("Trigger", "ユーザーが自発的に開く（プル型）"),
]
y = 2.5
for label, desc in flow_items:
    add_text(slide, 1.0, y, 1.5, 0.3, label, font_size=11, color=BLUE, bold=True)
    add_text(slide, 2.5, y, 3.5, 0.3, desc, font_size=11, color=GRAY)
    y += 0.35

add_text(slide, 0.8, 3.7, 5.4, 0.3, "= 家計の「ダッシュボード」（見に行く人だけに役立つ）",
         font_size=12, color=GRAY, bold=True)

# Orbit Budget
add_shape_rect(slide, 6.9, 1.4, 5.8, 2.8, NAVY_LIGHT, GREEN)
add_text(slide, 7.1, 1.5, 5.4, 0.4, "Orbit Budget（介入型）",
         font_size=16, color=GREEN, bold=True)
add_text(slide, 7.1, 2.0, 5.4, 0.3, "ロック画面を見る → 色が目に入る → 「やばい」→ 支出を止める",
         font_size=11, color=WHITE)
flow_items2 = [
    ("Motivation", "赤いウィジェットを緑に戻したい（内発的）"),
    ("Ability", "記録も閲覧も3秒。認知負荷ゼロ"),
    ("Trigger", "ロック画面で自動発火（プッシュ型）"),
]
y = 2.5
for label, desc in flow_items2:
    add_text(slide, 7.3, y, 1.5, 0.3, label, font_size=11, color=GREEN, bold=True)
    add_text(slide, 8.8, y, 3.5, 0.3, desc, font_size=11, color=WHITE)
    y += 0.35

add_text(slide, 7.1, 3.7, 5.4, 0.3, "= 家計の「警報装置」（見なくても鳴る）",
         font_size=12, color=GREEN, bold=True)

# 下部の分析
add_text(slide, 0.6, 4.5, 12, 0.4,
         "BJ Fogg の行動モデル（B = MAT）による比較",
         font_size=14, color=WHITE, bold=True)

add_text(slide, 0.6, 5.0, 12, 1.5,
         "キャッシュレス社会で消えた「財布の中身が減る感覚」を、"
         "ウィジェットの色変化で疑似的に復元する。\n"
         "これが Orbit Budget の本質的な価値 = 「デジタル時代の支払いの痛みの再注入」\n\n"
         "ただし、警告には順応効果がある。赤いウィジェットにも慣れる。\n"
         "→ チャレンジ機能・達成演出・日割り予算の変動表示で「慣れ」を防ぐ設計が必要",
         font_size=12, color=GRAY)

# ==========================================
# スライド5: ADHD当事者視点の厳正評価
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.6, 0.4, 12, 0.5, "ADHD当事者視点: 本当にこのアプリで浪費は止まるか？",
         font_size=28, color=WHITE, bold=True)
add_shape_rect(slide, 0.6, 1.0, 12.1, 0.03, GREEN)

# ADHD 4原因
causes = [
    ("衝動性", "「欲しい」→ 即購入。\nブレーキが効かない", ORANGE),
    ("実行機能障害", "「管理しよう」→ 行動に\n落とし込めない", BLUE),
    ("時間感覚の歪み", "月末まであとX日が\n認識できない", RED),
    ("報酬系の問題", "買い物＝手軽な\nドーパミン供給源", RGBColor(168, 85, 247)),
]
for i, (name, desc, color) in enumerate(causes):
    x = 0.6 + i * 3.1
    add_shape_rect(slide, x, 1.3, 2.9, 1.3, NAVY_LIGHT, color)
    add_text(slide, x + 0.15, 1.4, 2.6, 0.3, name,
             font_size=14, color=color, bold=True)
    add_text(slide, x + 0.15, 1.75, 2.6, 0.7, desc,
             font_size=11, color=GRAY)

# 機能評価
add_text(slide, 0.6, 2.85, 12, 0.4, "各機能の ADHD 適合度評価",
         font_size=16, color=WHITE, bold=True)

features = [
    ("ウィジェット常時表示", "○", "正しいが「見る」だけでは\n行動は変わらない"),
    ("3秒入力", "◎", "最大の武器。ただし\n入力自体を忘れる問題"),
    ("予算ゲージ", "○", "色変化は直感的だが\n情報量が多すぎる"),
    ("立替分離", "◎", "認知負荷の分離が\n完璧に機能"),
    ("チュートリアル", "△", "ADHDの人は\nスキップする"),
    ("通知", "×", "未実装。\nこれが致命的"),
]

for i, (name, rating, desc) in enumerate(features):
    x = 0.6 + (i % 3) * 4.2
    y = 3.3 + (i // 3) * 1.3
    rating_color = GREEN if rating in ("◎",) else (
        RGBColor(74, 222, 128) if rating == "○" else (ORANGE if rating == "△" else RED))
    add_shape_rect(slide, x, y, 3.9, 1.1, NAVY_LIGHT, RGBColor(40, 40, 60))
    add_text(slide, x + 0.15, y + 0.1, 2.5, 0.3, name,
             font_size=12, color=WHITE, bold=True)
    add_text(slide, x + 3.2, y + 0.1, 0.5, 0.3, rating,
             font_size=18, color=rating_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + 0.15, y + 0.45, 3.5, 0.6, desc,
             font_size=10, color=GRAY)

add_text(slide, 0.6, 6.2, 12, 0.8,
         "致命的な欠如: ① 購入「前」の介入がない ② 日割り予算がない ③ プッシュ通知がない\n"
         "→ 現在は「記録を楽にするアプリ」。「浪費を止めるアプリ」になるには上記3つが不可欠",
         font_size=12, color=ORANGE, bold=True)

# ==========================================
# スライド6: 3つの改善提案
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.6, 0.4, 12, 0.5, "「記録型」→「介入型」への転換: 3つの改善提案",
         font_size=28, color=WHITE, bold=True)
add_shape_rect(slide, 0.6, 1.0, 12.1, 0.03, GREEN)

proposals = [
    ("1", "「今日使える金額」の導入",
     "月の残り予算 ÷ 残り日数 = 今日の予算",
     "Time Blindness（時間感覚の歪み）への直接対策。\n"
     "「残り82,000円」より「今日は2,700円まで」の方が\n"
     "行動を変える力が圧倒的に強い。",
     "S（計算ロジック + 表示変更のみ）",
     "ウィジェット・ダッシュボードの\n最上位に表示",
     GREEN),
    ("2", "「購入前チェック」モード",
     "買う前に金額入力 → 残額シミュレーション → 確認",
     "衝動性に対する唯一の有効な対策は\n「意思決定と実行の間に物理的な間を作る」こと。\n"
     "この1ステップが最後のブレーキになる。",
     "M（新モード UI + フロー追加）",
     "Quick Input に「これから使う」\nモードを追加",
     ORANGE),
    ("3", "記録忘れ防止リマインダー",
     "毎日21時に通知 + 未記録日の検知",
     "ADHDが家計簿を3日で放棄する最大の原因は\n「空白期間ができてやる気を失う」こと。\n"
     "ポジティブな文言で習慣の外部足場を提供。",
     "S（ローカル通知のみ）",
     "UNUserNotificationCenter で\nローカル通知を実装",
     BLUE),
]

for i, (num, title, how, why, effort, impl, color) in enumerate(proposals):
    x = 0.6 + i * 4.1
    add_shape_rect(slide, x, 1.3, 3.9, 5.5, NAVY_LIGHT, color)

    # 番号
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(x + 0.15), Inches(1.45), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text(slide, x + 0.15, 1.45, 0.45, 0.45, num,
             font_size=20, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

    add_text(slide, x + 0.7, 1.48, 3.0, 0.4, title,
             font_size=16, color=WHITE, bold=True)

    add_text(slide, x + 0.2, 2.05, 3.5, 0.3, "What:", font_size=10, color=color, bold=True)
    add_text(slide, x + 0.2, 2.3, 3.5, 0.5, how, font_size=11, color=WHITE)

    add_text(slide, x + 0.2, 2.9, 3.5, 0.3, "Why:", font_size=10, color=color, bold=True)
    add_text(slide, x + 0.2, 3.15, 3.5, 1.2, why, font_size=10, color=GRAY)

    add_text(slide, x + 0.2, 4.5, 3.5, 0.3, "実装:", font_size=10, color=color, bold=True)
    add_text(slide, x + 0.2, 4.75, 3.5, 0.5, impl, font_size=10, color=GRAY)

    add_text(slide, x + 0.2, 5.5, 3.5, 0.3, "工数:", font_size=10, color=color, bold=True)
    add_text(slide, x + 0.2, 5.75, 3.5, 0.3, effort, font_size=12, color=WHITE, bold=True)

add_text(slide, 0.6, 7.0, 12, 0.4,
         "3つとも現在のアーキテクチャに大きな変更なく追加可能。合計工数: 2〜3スプリント",
         font_size=12, color=GREEN, bold=True)

# ==========================================
# スライド7: ビジネスモデルの構造課題
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.6, 0.4, 12, 0.5, "ビジネスモデルの構造課題と対策",
         font_size=28, color=WHITE, bold=True)
add_shape_rect(slide, 0.6, 1.0, 12.1, 0.03, GREEN)

# 矛盾の図
add_text(slide, 0.6, 1.3, 12, 0.4, "構造的矛盾: プロダクトの成功 = ユーザーの離脱",
         font_size=18, color=RED, bold=True)

steps = [
    "浪費が止まる（成功）",
    "赤字にならない",
    "警告が出ない",
    "課金動機が消える",
    "LTVが下がる"
]
for i, step in enumerate(steps):
    x = 0.8 + i * 2.4
    c = GREEN if i == 0 else (ORANGE if i < 3 else RED)
    add_shape_rect(slide, x, 1.85, 2.1, 0.5, NAVY_LIGHT, c)
    add_text(slide, x + 0.1, 1.9, 1.9, 0.4, step,
             font_size=11, color=c, bold=True, alignment=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text(slide, x + 2.15, 1.9, 0.3, 0.4, "→",
                 font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# 対策
add_text(slide, 0.6, 2.7, 12, 0.4, "対策: ゴールの拡張",
         font_size=16, color=GREEN, bold=True)

add_shape_rect(slide, 0.6, 3.2, 5.8, 1.5, NAVY_LIGHT, BLUE)
add_text(slide, 0.8, 3.3, 5.4, 0.3, "現在のゴール（Phase 1-2）",
         font_size=13, color=BLUE, bold=True)
add_text(slide, 0.8, 3.65, 5.4, 0.8,
         "「浪費を止める」= 予算内で生活する\n\n"
         "→ 達成するとアプリ不要に",
         font_size=12, color=GRAY)

add_shape_rect(slide, 6.9, 3.2, 5.8, 1.5, NAVY_LIGHT, GREEN)
add_text(slide, 7.1, 3.3, 5.4, 0.3, "拡張ゴール（Phase 3以降）",
         font_size=13, color=GREEN, bold=True)
add_text(slide, 7.1, 3.65, 5.4, 0.8,
         "「お金と上手に付き合う」= 節約 → 貯蓄 → 投資\n\n"
         "→ 浪費が止まった後も使い続ける理由がある",
         font_size=12, color=WHITE)

# 市場規模
add_text(slide, 0.6, 5.0, 12, 0.4, "市場規模と現実的な収益見通し",
         font_size=16, color=WHITE, bold=True)

metrics = [
    ("TAM", "家計簿アプリ市場", "3,000万人", BLUE),
    ("SAM", "浪費家セグメント", "150〜250万人", GREEN),
    ("SOM", "初年度獲得目標", "1〜3万人", ORANGE),
    ("黒字化ライン", "広告+サブスク", "DAU 3,000人", RED),
]
for i, (label, desc, value, color) in enumerate(metrics):
    x = 0.6 + i * 3.1
    add_shape_rect(slide, x, 5.5, 2.9, 1.3, NAVY_LIGHT, color)
    add_text(slide, x + 0.15, 5.6, 2.6, 0.25, label,
             font_size=11, color=color, bold=True)
    add_text(slide, x + 0.15, 5.85, 2.6, 0.3, value,
             font_size=18, color=WHITE, bold=True)
    add_text(slide, x + 0.15, 6.25, 2.6, 0.3, desc,
             font_size=10, color=GRAY)

# ==========================================
# スライド8: ターゲットセグメント
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.6, 0.4, 12, 0.5, "ターゲットセグメント: 誰を最初に攻めるか",
         font_size=28, color=WHITE, bold=True)
add_shape_rect(slide, 0.6, 1.0, 12.1, 0.03, GREEN)

segments = [
    ("B: 飲み会社畜型", "★★★ 最優先", "25〜32歳 / 手取り22〜35万", "飲み会・外食・ガジェット",
     "立替セパレーターが直接刺さる\n既存アプリへの不満が具体的", GREEN),
    ("E: ADHD当事者", "★★★ 最優先", "全年齢 / 全帯域", "衝動買い全般・サブスク忘れ",
     "「ADHD向け」がそのまま刺さる\n自助コミュニティで口コミ拡散", GREEN),
    ("A: 衝動買いSNS世代", "★★☆ 第2優先", "20〜25歳 / 手取り15〜22万", "コスメ・カフェ・推し活",
     "WTPが低い（月480円への抵抗）\n推し活は「浪費」と認識していない", ORANGE),
    ("C: 見栄消費カップル", "★☆☆ 第3優先", "28〜35歳 / 手取り25〜40万", "デート・旅行・ブランド",
     "個人アプリでは対応しにくい\n家族向け機能が必要", GRAY),
]

for i, (name, priority, demo, pattern, reason, color) in enumerate(segments):
    y = 1.3 + i * 1.45
    add_shape_rect(slide, 0.6, y, 12.1, 1.3, NAVY_LIGHT, color)
    add_text(slide, 0.8, y + 0.1, 3.0, 0.3, name,
             font_size=16, color=color, bold=True)
    add_text(slide, 3.8, y + 0.1, 2.0, 0.3, priority,
             font_size=12, color=color, bold=True)
    add_text(slide, 0.8, y + 0.5, 2.5, 0.3, demo,
             font_size=11, color=GRAY)
    add_text(slide, 3.5, y + 0.5, 3.0, 0.3, pattern,
             font_size=11, color=GRAY)
    add_text(slide, 7.0, y + 0.15, 5.5, 0.8, reason,
             font_size=11, color=WHITE)

add_text(slide, 0.6, 7.0, 12, 0.4,
         "戦略: セグメントBとEの二刀流でPMFを検証 → D30継続率20%達成後にセグメントAへ拡大",
         font_size=12, color=GREEN, bold=True)

# ==========================================
# スライド9: Next Steps
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.6, 0.4, 12, 0.5, "NEXT STEPS: 今後のアクションプラン",
         font_size=28, color=WHITE, bold=True)
add_shape_rect(slide, 0.6, 1.0, 12.1, 0.03, GREEN)

phases = [
    ("NOW", "MVP リリース準備", [
        "TestFlight 配布（現在のコードベースで）",
        "ST（システムテスト）10件の実施",
        "App Store 提出・審査",
    ], GREEN),
    ("v1.1", "「介入型」機能追加", [
        "①「今日使える金額」の実装（工数 S）",
        "③ 記録リマインダー通知の実装（工数 S）",
        "ウィジェットへの日割り予算表示",
    ], BLUE),
    ("v1.2", "行動変容の深化", [
        "② 購入前チェックモードの実装（工数 M）",
        "チャレンジ機能・ストリーク",
        "週間/月間レポート通知",
    ], ORANGE),
    ("v2.0", "ゴールの拡張", [
        "「お金と上手に付き合う」へシフト",
        "貯蓄目標機能",
        "AdMob + サブスク本格導入",
    ], RGBColor(168, 85, 247)),
]

for i, (version, title, items, color) in enumerate(phases):
    x = 0.6 + i * 3.1
    add_shape_rect(slide, x, 1.3, 2.9, 4.5, NAVY_LIGHT, color)

    add_text(slide, x + 0.15, 1.4, 2.6, 0.3, version,
             font_size=22, color=color, bold=True)
    add_text(slide, x + 0.15, 1.8, 2.6, 0.3, title,
             font_size=14, color=WHITE, bold=True)

    y = 2.3
    for item in items:
        add_text(slide, x + 0.15, y, 2.6, 0.5, f"• {item}",
                 font_size=11, color=GRAY)
        y += 0.45

    if i < len(phases) - 1:
        add_text(slide, x + 2.95, 3.0, 0.3, 0.5, "→",
                 font_size=24, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text(slide, 0.6, 6.2, 12, 0.8,
         "最重要KPI: D30継続率 20% の達成\n"
         "マネタイズは PMF（Product-Market Fit）のシグナルが確認できてからでも遅くない",
         font_size=14, color=GREEN, bold=True)

# 保存
output_path = "/Users/taguchikoji/Desktop/vibecoding/second-business/No-Look-Budget/docs/reports/orbit_budget_strategic_review.pptx"
prs.save(output_path)
print(f"PowerPoint saved: {output_path}")
